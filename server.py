#!/usr/bin/env python3
"""
SERVIDOR DE CLASES DE INGLÉS - VERSIÓN 7.0
Con DeepSeek streaming, búsqueda en internet y procesamiento de archivos
Solo PostgreSQL (sin SQLite)
"""

import os
import sys
import json
import hashlib
import secrets
import jwt
import datetime
import base64
import io
import time
from typing import Dict, Optional, Generator
from functools import wraps
from urllib.parse import urlparse

from flask import Flask, request, jsonify, send_file, send_from_directory, Response, stream_with_context
from flask_cors import CORS
from werkzeug.utils import secure_filename
from flask_socketio import SocketIO, emit

# PostgreSQL
import psycopg2
from psycopg2.extras import RealDictCursor

# Cloudinary
import cloudinary
import cloudinary.uploader

# DeepSeek / OpenAI
import openai

# Procesamiento de archivos
import PyPDF2
from docx import Document
from pptx import Presentation
import openpyxl
from PIL import Image
import requests
from io import BytesIO

# ============ CONFIGURACIÓN DIRECTA ============
app = Flask(__name__, static_folder='.', static_url_path='')
CORS(app, resources={r"/*": {"origins": "*"}})

app.config['SECRET_KEY'] = 'clases-ingles-secret-key-2024'
app.config['MAX_CONTENT_LENGTH'] = 50 * 1024 * 1024  # 50MB

# ============ CONFIGURACIÓN POSTGRESQL ============
DATABASE_URL = "postgresql://englishcourse_user:VI8pYTtX2bbv2YftidVHOUKXtK6J7ehd@dpg-d64mmungi27c73b53hr0-a/englishcourse"

# ============ CONFIGURACIÓN CLOUDINARY ============
cloudinary.config(
    cloud_name="dj72b0ykc",
    api_key="215156196366932",
    api_secret="Ivdpe_mkT3rSx5asFTo6qJdWaLQ",
    secure=True
)

# ============ DEEPSEEK CONFIGURACIÓN ============
DEEPSEEK_API_KEY = "sk-0c21275ee63f4b42ac0733835ac44d29"

# Cliente de DeepSeek (compatible con OpenAI)
deepseek_client = openai.OpenAI(
    api_key=DEEPSEEK_API_KEY,
    base_url="https://api.deepseek.com"
)

# ============ CÓDIGOS DE ACCESO ============
STUDENT_CODE = "QwErTy89"
TEACHER_CODE = "MOISÉS5M"

# Extensiones permitidas
ALLOWED_EXTENSIONS = {
    'pdf', 'doc', 'docx', 'ppt', 'pptx', 'xls', 'xlsx',
    'txt', 'jpg', 'jpeg', 'png', 'gif',
    'mp3', 'mp4', 'mov', 'avi', 'webm'
}

socketio = SocketIO(app, cors_allowed_origins="*", ping_timeout=60, ping_interval=25)

# ============ FUNCIONES AUXILIARES ============
def allowed_file(filename: str) -> bool:
    return '.' in filename and \
           filename.rsplit('.', 1)[1].lower() in ALLOWED_EXTENSIONS

def get_file_type(filename: str) -> str:
    ext = filename.rsplit('.', 1)[1].lower() if '.' in filename else ''
    file_types = {
        'pdf': 'PDF', 'doc': 'Word', 'docx': 'Word',
        'ppt': 'PowerPoint', 'pptx': 'PowerPoint',
        'xls': 'Excel', 'xlsx': 'Excel',
        'txt': 'Texto', 'jpg': 'Imagen', 'jpeg': 'Imagen',
        'png': 'Imagen', 'gif': 'GIF', 'mp3': 'Audio',
        'mp4': 'Video', 'mov': 'Video', 'avi': 'Video',
        'webm': 'Video'
    }
    return file_types.get(ext, 'Archivo')

def extract_text_from_file(file, filename: str) -> str:
    """Extrae texto de diferentes tipos de archivos"""
    ext = filename.rsplit('.', 1)[1].lower() if '.' in filename else ''
    
    try:
        if ext == 'txt':
            return file.read().decode('utf-8', errors='ignore')
        
        elif ext == 'pdf':
            pdf_reader = PyPDF2.PdfReader(file)
            text = ""
            for page in pdf_reader.pages:
                text += page.extract_text()
            return text
        
        elif ext in ['docx']:
            doc = Document(io.BytesIO(file.read()))
            return "\n".join([paragraph.text for paragraph in doc.paragraphs])
        
        elif ext in ['pptx']:
            prs = Presentation(io.BytesIO(file.read()))
            text = ""
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
            return text
        
        elif ext in ['xlsx']:
            wb = openpyxl.load_workbook(io.BytesIO(file.read()), data_only=True)
            text = ""
            for sheet in wb.sheetnames:
                ws = wb[sheet]
                for row in ws.iter_rows(values_only=True):
                    text += " ".join([str(cell) for cell in row if cell]) + "\n"
            return text
        
        else:
            return f"[Archivo {ext} no procesable para texto]"
            
    except Exception as e:
        return f"[Error extrayendo texto: {str(e)}]"

def get_db_connection():
    """Obtiene una conexión a PostgreSQL"""
    return psycopg2.connect(DATABASE_URL, cursor_factory=RealDictCursor)

def init_db():
    """Inicializar base de datos PostgreSQL"""
    conn = get_db_connection()
    cursor = conn.cursor()
    
    try:
        # Tabla de usuarios
        cursor.execute('''
            CREATE TABLE IF NOT EXISTS users (
                id SERIAL PRIMARY KEY,
                username VARCHAR(100) UNIQUE NOT NULL,
                password_hash VARCHAR(255) NOT NULL,
                name VARCHAR(200) NOT NULL,
                role VARCHAR(20) NOT NULL CHECK(role IN ('student', 'teacher')),
                level VARCHAR(50) DEFAULT 'Sin asignar',
                avatar_url TEXT,
                avatar_version INTEGER DEFAULT 1,
                registration_date TIMESTAMP DEFAULT CURRENT_TIMESTAMP,
                last_login TIMESTAMP,
                is_active BOOLEAN DEFAULT TRUE
            )
        ''')
        
        # Tabla de clases
        cursor.execute('''
            CREATE TABLE IF NOT EXISTS classes (
                id SERIAL PRIMARY KEY,
                title VARCHAR(255) NOT NULL,
                description TEXT NOT NULL,
                level VARCHAR(10) NOT NULL CHECK(level IN ('A1', 'A2', 'B1')),
                file_name VARCHAR(255),
                file_url TEXT,
                file_type VARCHAR(50),
                file_size INTEGER,
                uploaded_by INTEGER NOT NULL,
                created_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP,
                is_active BOOLEAN DEFAULT TRUE,
                FOREIGN KEY (uploaded_by) REFERENCES users (id) ON DELETE CASCADE
            )
        ''')
        
        # Tabla de acceso estudiante-clase
        cursor.execute('''
            CREATE TABLE IF NOT EXISTS student_access (
                id SERIAL PRIMARY KEY,
                student_id INTEGER NOT NULL,
                class_id INTEGER NOT NULL,
                has_access BOOLEAN DEFAULT FALSE,
                downloaded BOOLEAN DEFAULT FALSE,
                downloaded_at TIMESTAMP,
                granted_by INTEGER,
                granted_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP,
                FOREIGN KEY (student_id) REFERENCES users (id) ON DELETE CASCADE,
                FOREIGN KEY (class_id) REFERENCES classes (id) ON DELETE CASCADE,
                FOREIGN KEY (granted_by) REFERENCES users (id) ON DELETE SET NULL,
                UNIQUE(student_id, class_id)
            )
        ''')
        
        # Tabla de descargas
        cursor.execute('''
            CREATE TABLE IF NOT EXISTS downloads (
                id SERIAL PRIMARY KEY,
                user_id INTEGER NOT NULL,
                class_id INTEGER NOT NULL,
                downloaded_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP,
                FOREIGN KEY (user_id) REFERENCES users (id) ON DELETE CASCADE,
                FOREIGN KEY (class_id) REFERENCES classes (id) ON DELETE CASCADE
            )
        ''')
        
        # Tabla de publicaciones
        cursor.execute('''
            CREATE TABLE IF NOT EXISTS publications (
                id SERIAL PRIMARY KEY,
                title VARCHAR(255) NOT NULL,
                content TEXT,
                publication_type VARCHAR(20) NOT NULL CHECK(publication_type IN ('text', 'photo', 'video', 'document')),
                file_url TEXT,
                file_name VARCHAR(255),
                file_type VARCHAR(50),
                created_by INTEGER NOT NULL,
                created_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP,
                is_active BOOLEAN DEFAULT TRUE,
                FOREIGN KEY (created_by) REFERENCES users (id) ON DELETE CASCADE
            )
        ''')
        
        # Tabla de likes
        cursor.execute('''
            CREATE TABLE IF NOT EXISTS publication_likes (
                id SERIAL PRIMARY KEY,
                user_id INTEGER NOT NULL,
                publication_id INTEGER NOT NULL,
                liked_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP,
                FOREIGN KEY (user_id) REFERENCES users (id) ON DELETE CASCADE,
                FOREIGN KEY (publication_id) REFERENCES publications (id) ON DELETE CASCADE,
                UNIQUE(user_id, publication_id)
            )
        ''')
        
        conn.commit()
        print("✅ Base de datos PostgreSQL configurada")
        
        # Crear profesor por defecto si no existe
        cursor.execute("SELECT id FROM users WHERE role = 'teacher' LIMIT 1")
        if not cursor.fetchone():
            password_hash = hash_password("admin123")
            cursor.execute('''
                INSERT INTO users (username, password_hash, name, role)
                VALUES (%s, %s, %s, %s)
            ''', ('profesor', password_hash, 'Profesor Moisés', 'teacher'))
            conn.commit()
            print("✅ Profesor por defecto: usuario='profesor', password='admin123'")
        
    except Exception as e:
        print(f"⚠️ Error: {str(e)}")
        conn.rollback()
    finally:
        cursor.close()
        conn.close()

def hash_password(password: str) -> str:
    return hashlib.sha256(password.encode()).hexdigest()

def verify_password(password: str, password_hash: str) -> bool:
    return hash_password(password) == password_hash

def generate_token(user_id: int, username: str, role: str) -> str:
    payload = {
        'user_id': user_id,
        'username': username,
        'role': role,
        'exp': datetime.datetime.utcnow() + datetime.timedelta(days=7)
    }
    return jwt.encode(payload, app.config['SECRET_KEY'], algorithm='HS256')

def verify_token(token: str) -> Optional[Dict]:
    try:
        return jwt.decode(token, app.config['SECRET_KEY'], algorithms=['HS256'])
    except:
        return None

# ============ DECORADORES ============
def token_required(f):
    @wraps(f)
    def decorated(*args, **kwargs):
        token = None
        auth_header = request.headers.get('Authorization')
        
        if auth_header and auth_header.startswith('Bearer '):
            token = auth_header.split(' ')[1]
        
        if not token:
            return jsonify({'error': 'Token requerido'}), 401
        
        payload = verify_token(token)
        if not payload:
            return jsonify({'error': 'Token inválido'}), 401
        
        request.user_id = payload['user_id']
        request.username = payload['username']
        request.user_role = payload['role']
        
        return f(*args, **kwargs)
    return decorated

def teacher_required(f):
    @wraps(f)
    def decorated(*args, **kwargs):
        if not hasattr(request, 'user_role') or request.user_role != 'teacher':
            return jsonify({'error': 'Acceso denegado'}), 403
        return f(*args, **kwargs)
    return decorated

# ============ RUTAS PRINCIPALES ============
@app.route('/')
def index():
    return send_file('auth.html')

@app.route('/auth.html')
def auth_page():
    return send_file('auth.html')

@app.route('/estudiante.html')
def estudiante_page():
    return send_file('estudiante.html')

@app.route('/profesor.html')
def profesor_page():
    return send_file('profesor.html')

# ============ AUTENTICACIÓN ============
@app.route('/api/register', methods=['POST'])
def register():
    try:
        data = request.get_json()
        
        if not data:
            return jsonify({'error': 'Datos inválidos'}), 400
        
        required_fields = ['name', 'username', 'password', 'code']
        for field in required_fields:
            if not data.get(field):
                return jsonify({'error': f'Campo {field} requerido'}), 400
        
        code = data['code']
        if code not in [STUDENT_CODE, TEACHER_CODE]:
            return jsonify({'error': 'Código incorrecto'}), 400
        
        role = 'teacher' if code == TEACHER_CODE else 'student'
        
        conn = get_db_connection()
        cursor = conn.cursor()
        
        cursor.execute('SELECT id FROM users WHERE username = %s', (data['username'],))
        if cursor.fetchone():
            cursor.close()
            conn.close()
            return jsonify({'error': 'Usuario ya existe'}), 400
        
        password_hash = hash_password(data['password'])
        
        cursor.execute('''
            INSERT INTO users (username, password_hash, name, role, level)
            VALUES (%s, %s, %s, %s, %s) RETURNING id, username, name, role, level, registration_date, avatar_url, avatar_version
        ''', (
            data['username'],
            password_hash,
            data['name'],
            role,
            'Sin asignar' if role == 'student' else None
        ))
        
        user = cursor.fetchone()
        conn.commit()
        
        token = generate_token(user['id'], user['username'], user['role'])
        
        cursor.close()
        conn.close()
        
        return jsonify({
            'message': 'Registro exitoso',
            'user': dict(user),
            'token': token
        }), 201
        
    except Exception as e:
        print(f"❌ Error: {str(e)}")
        return jsonify({'error': 'Error interno'}), 500

@app.route('/api/login', methods=['POST'])
def login():
    try:
        data = request.get_json()
        
        if not data or not data.get('username') or not data.get('password'):
            return jsonify({'error': 'Usuario y contraseña requeridos'}), 400
        
        conn = get_db_connection()
        cursor = conn.cursor()
        
        cursor.execute('SELECT * FROM users WHERE username = %s', (data['username'],))
        user_row = cursor.fetchone()
        
        if not user_row:
            cursor.close()
            conn.close()
            return jsonify({'error': 'Credenciales incorrectas'}), 401
        
        user = dict(user_row)
        
        if not verify_password(data['password'], user['password_hash']):
            cursor.close()
            conn.close()
            return jsonify({'error': 'Credenciales incorrectas'}), 401
        
        cursor.execute('UPDATE users SET last_login = CURRENT_TIMESTAMP WHERE id = %s', (user['id'],))
        conn.commit()
        
        del user['password_hash']
        
        token = generate_token(user['id'], user['username'], user['role'])
        
        cursor.close()
        conn.close()
        
        redirect_url = 'profesor.html' if user['role'] == 'teacher' else 'estudiante.html'
        
        return jsonify({
            'message': 'Login exitoso',
            'user': user,
            'token': token,
            'redirect': redirect_url
        }), 200
        
    except Exception as e:
        print(f"❌ Error: {str(e)}")
        return jsonify({'error': 'Error interno'}), 500

# ============ AVATAR ============
@app.route('/api/profile/avatar', methods=['POST'])
@token_required
def upload_avatar():
    try:
        if 'avatar' not in request.files:
            return jsonify({'error': 'No se envió archivo'}), 400
        
        file = request.files['avatar']
        
        if file.filename == '':
            return jsonify({'error': 'Archivo vacío'}), 400
        
        if not allowed_file(file.filename):
            return jsonify({'error': 'Formato no permitido'}), 400
        
        upload_result = cloudinary.uploader.upload(
            file,
            folder="avatars",
            resource_type="image",
            use_filename=True,
            unique_filename=True,
            width=200,
            height=200,
            crop="fill"
        )
        
        conn = get_db_connection()
        cursor = conn.cursor()
        
        cursor.execute('''
            UPDATE users 
            SET avatar_url = %s, avatar_version = avatar_version + 1
            WHERE id = %s
            RETURNING avatar_url, avatar_version
        ''', (upload_result['secure_url'], request.user_id))
        
        result = cursor.fetchone()
        conn.commit()
        cursor.close()
        conn.close()
        
        return jsonify({
            'message': 'Avatar actualizado',
            'avatar_url': result['avatar_url'],
            'avatar_version': result['avatar_version']
        }), 200
        
    except Exception as e:
        print(f"❌ Error: {str(e)}")
        return jsonify({'error': 'Error al subir avatar'}), 500

# ============ PUBLICACIONES ============
@app.route('/api/publications', methods=['GET'])
@token_required
def get_publications():
    try:
        conn = get_db_connection()
        cursor = conn.cursor()
        
        cursor.execute('''
            SELECT p.*, u.name as author_name, u.avatar_url, u.avatar_version,
                   (SELECT COUNT(*) FROM publication_likes pl WHERE pl.publication_id = p.id) as likes_count,
                   EXISTS(SELECT 1 FROM publication_likes pl WHERE pl.publication_id = p.id AND pl.user_id = %s) as user_liked
            FROM publications p
            JOIN users u ON p.created_by = u.id
            WHERE p.is_active = TRUE
            ORDER BY p.created_at DESC
        ''', (request.user_id,))
        
        publications = []
        for row in cursor.fetchall():
            publications.append(dict(row))
        
        cursor.close()
        conn.close()
        
        return jsonify(publications), 200
        
    except Exception as e:
        print(f"❌ Error: {str(e)}")
        return jsonify({'error': 'Error obteniendo publicaciones'}), 500

@app.route('/api/publications/<int:publication_id>/like', methods=['POST'])
@token_required
def toggle_like(publication_id: int):
    try:
        if request.user_role != 'student':
            return jsonify({'error': 'Solo estudiantes pueden dar like'}), 403
        
        conn = get_db_connection()
        cursor = conn.cursor()
        
        cursor.execute('''
            SELECT id FROM publication_likes 
            WHERE user_id = %s AND publication_id = %s
        ''', (request.user_id, publication_id))
        
        existing = cursor.fetchone()
        
        if existing:
            cursor.execute('''
                DELETE FROM publication_likes 
                WHERE user_id = %s AND publication_id = %s
            ''', (request.user_id, publication_id))
            liked = False
        else:
            cursor.execute('''
                INSERT INTO publication_likes (user_id, publication_id)
                VALUES (%s, %s)
            ''', (request.user_id, publication_id))
            liked = True
        
        conn.commit()
        
        cursor.execute('''
            SELECT COUNT(*) as count FROM publication_likes 
            WHERE publication_id = %s
        ''', (publication_id,))
        count = cursor.fetchone()['count']
        
        cursor.close()
        conn.close()
        
        socketio.emit('like_updated', {
            'publication_id': publication_id,
            'likes_count': count,
            'user_id': request.user_id,
            'liked': liked
        })
        
        return jsonify({
            'liked': liked,
            'likes_count': count
        }), 200
        
    except Exception as e:
        print(f"❌ Error: {str(e)}")
        return jsonify({'error': 'Error'}), 500

@app.route('/api/publications', methods=['POST'])
@token_required
@teacher_required
def create_publication():
    try:
        publication_type = request.form.get('type', 'text')
        
        if publication_type not in ['text', 'photo', 'video', 'document']:
            return jsonify({'error': 'Tipo no válido'}), 400
        
        title = request.form.get('title', '').strip()
        content = request.form.get('content', '').strip()
        
        if not title:
            return jsonify({'error': 'Título requerido'}), 400
        
        if publication_type == 'text' and not content:
            return jsonify({'error': 'Contenido requerido'}), 400
        
        file_url = None
        file_name = None
        file_type = None
        
        if publication_type in ['photo', 'video', 'document']:
            if 'file' not in request.files:
                return jsonify({'error': 'Archivo requerido'}), 400
            
            file = request.files['file']
            
            if file.filename == '':
                return jsonify({'error': 'Archivo vacío'}), 400
            
            if not allowed_file(file.filename):
                return jsonify({'error': 'Tipo no permitido'}), 400
            
            resource_type = 'auto'
            if publication_type == 'video':
                resource_type = 'video'
            elif publication_type == 'photo':
                resource_type = 'image'
            
            upload_result = cloudinary.uploader.upload(
                file,
                folder="publicaciones",
                resource_type=resource_type,
                use_filename=True,
                unique_filename=True
            )
            
            file_url = upload_result['secure_url']
            file_name = secure_filename(file.filename)
            file_type = get_file_type(file.filename)
        
        conn = get_db_connection()
        cursor = conn.cursor()
        
        cursor.execute('SELECT name, avatar_url, avatar_version FROM users WHERE id = %s', (request.user_id,))
        teacher = cursor.fetchone()
        
        cursor.execute('''
            INSERT INTO publications 
            (title, content, publication_type, file_url, file_name, file_type, created_by)
            VALUES (%s, %s, %s, %s, %s, %s, %s)
            RETURNING id, title, created_at
        ''', (
            title,
            content if publication_type == 'text' else '',
            publication_type,
            file_url,
            file_name,
            file_type,
            request.user_id
        ))
        
        new_publication = cursor.fetchone()
        conn.commit()
        cursor.close()
        conn.close()
        
        socketio.emit('new_publication', {
            'id': new_publication['id'],
            'title': title,
            'author_name': teacher['name'],
            'avatar_url': teacher['avatar_url'],
            'avatar_version': teacher['avatar_version'],
            'created_at': new_publication['created_at'].isoformat()
        })
        
        return jsonify({
            'message': 'Publicación creada',
            'publication': dict(new_publication)
        }), 201
        
    except Exception as e:
        print(f"❌ Error: {str(e)}")
        return jsonify({'error': 'Error'}), 500

@app.route('/api/publications/<int:publication_id>', methods=['DELETE'])
@token_required
@teacher_required
def delete_publication(publication_id: int):
    try:
        conn = get_db_connection()
        cursor = conn.cursor()
        
        cursor.execute('UPDATE publications SET is_active = FALSE WHERE id = %s', (publication_id,))
        conn.commit()
        
        cursor.close()
        conn.close()
        
        socketio.emit('publication_deleted', {'id': publication_id})
        
        return jsonify({'message': 'Publicación eliminada'}), 200
        
    except Exception as e:
        print(f"❌ Error: {str(e)}")
        return jsonify({'error': 'Error'}), 500

# ============ CLASES ============
@app.route('/api/classes', methods=['GET'])
@token_required
def get_classes():
    try:
        level = request.args.get('level')
        
        conn = get_db_connection()
        cursor = conn.cursor()
        
        if request.user_role == 'teacher':
            if level:
                cursor.execute('''
                    SELECT c.*, u.name as uploaded_by_name 
                    FROM classes c 
                    JOIN users u ON c.uploaded_by = u.id 
                    WHERE c.level = %s AND c.is_active = TRUE
                    ORDER BY c.created_at DESC
                ''', (level,))
            else:
                cursor.execute('''
                    SELECT c.*, u.name as uploaded_by_name 
                    FROM classes c 
                    JOIN users u ON c.uploaded_by = u.id 
                    WHERE c.is_active = TRUE
                    ORDER BY c.created_at DESC
                ''')
        else:
            cursor.execute('''
                SELECT c.*, u.name as uploaded_by_name, 
                       sa.has_access, sa.downloaded
                FROM classes c 
                JOIN users u ON c.uploaded_by = u.id 
                LEFT JOIN student_access sa ON c.id = sa.class_id AND sa.student_id = %s
                WHERE c.level = %s AND c.is_active = TRUE AND sa.has_access = TRUE
                ORDER BY c.created_at DESC
            ''', (request.user_id, level))
        
        classes = []
        for row in cursor.fetchall():
            classes.append(dict(row))
        
        cursor.close()
        conn.close()
        
        return jsonify(classes), 200
        
    except Exception as e:
        print(f"❌ Error: {str(e)}")
        return jsonify({'error': 'Error'}), 500

@app.route('/api/classes/all', methods=['GET'])
@token_required
@teacher_required
def get_all_classes():
    try:
        conn = get_db_connection()
        cursor = conn.cursor()
        
        cursor.execute('''
            SELECT c.*, u.name as uploaded_by_name 
            FROM classes c 
            JOIN users u ON c.uploaded_by = u.id 
            WHERE c.is_active = TRUE
            ORDER BY c.created_at DESC
        ''')
        
        classes = []
        for row in cursor.fetchall():
            classes.append(dict(row))
        
        cursor.close()
        conn.close()
        
        return jsonify(classes), 200
        
    except Exception as e:
        print(f"❌ Error: {str(e)}")
        return jsonify({'error': 'Error'}), 500

@app.route('/api/classes/<int:class_id>', methods=['GET'])
@token_required
def get_class(class_id: int):
    try:
        conn = get_db_connection()
        cursor = conn.cursor()
        
        if request.user_role == 'teacher':
            cursor.execute('''
                SELECT c.*, u.name as uploaded_by_name 
                FROM classes c 
                JOIN users u ON c.uploaded_by = u.id 
                WHERE c.id = %s
            ''', (class_id,))
        else:
            cursor.execute('''
                SELECT c.*, u.name as uploaded_by_name, sa.has_access
                FROM classes c 
                JOIN users u ON c.uploaded_by = u.id 
                LEFT JOIN student_access sa ON c.id = sa.class_id AND sa.student_id = %s
                WHERE c.id = %s AND sa.has_access = TRUE
            ''', (request.user_id, class_id))
        
        class_data = cursor.fetchone()
        
        if not class_data:
            cursor.close()
            conn.close()
            return jsonify({'error': 'Clase no encontrada'}), 404
        
        cursor.close()
        conn.close()
        
        return jsonify(dict(class_data)), 200
        
    except Exception as e:
        print(f"❌ Error: {str(e)}")
        return jsonify({'error': 'Error'}), 500

@app.route('/api/classes', methods=['POST'])
@token_required
@teacher_required
def upload_class():
    try:
        if 'file' not in request.files:
            return jsonify({'error': 'Archivo requerido'}), 400
        
        file = request.files['file']
        
        if file.filename == '':
            return jsonify({'error': 'Archivo vacío'}), 400
        
        if not allowed_file(file.filename):
            return jsonify({'error': 'Tipo no permitido'}), 400
        
        title = request.form.get('title', '').strip()
        description = request.form.get('description', '').strip()
        level = request.form.get('level', '').strip()
        
        if not title or not description or not level:
            return jsonify({'error': 'Todos los campos requeridos'}), 400
        
        if level not in ['A1', 'A2', 'B1']:
            return jsonify({'error': 'Nivel no válido'}), 400
        
        upload_result = cloudinary.uploader.upload(
            file,
            folder="clases_ingles",
            resource_type="auto",
            use_filename=True,
            unique_filename=True
        )
        
        conn = get_db_connection()
        cursor = conn.cursor()
        
        cursor.execute('''
            INSERT INTO classes 
            (title, description, level, file_name, file_url, file_type, file_size, uploaded_by)
            VALUES (%s, %s, %s, %s, %s, %s, %s, %s)
            RETURNING id, title, level, created_at
        ''', (
            title,
            description,
            level,
            secure_filename(file.filename),
            upload_result['secure_url'],
            get_file_type(file.filename),
            upload_result.get('bytes', 0),
            request.user_id
        ))
        
        new_class = cursor.fetchone()
        conn.commit()
        cursor.close()
        conn.close()
        
        socketio.emit('new_class', {
            'id': new_class['id'],
            'title': title,
            'level': level,
            'created_at': new_class['created_at'].isoformat()
        })
        
        return jsonify({
            'message': 'Clase subida',
            'class': dict(new_class)
        }), 201
        
    except Exception as e:
        print(f"❌ Error: {str(e)}")
        return jsonify({'error': 'Error'}), 500

@app.route('/api/classes/<int:class_id>', methods=['PUT'])
@token_required
@teacher_required
def update_class(class_id: int):
    try:
        conn = get_db_connection()
        cursor = conn.cursor()
        
        cursor.execute('SELECT * FROM classes WHERE id = %s', (class_id,))
        if not cursor.fetchone():
            cursor.close()
            conn.close()
            return jsonify({'error': 'Clase no encontrada'}), 404
        
        title = request.form.get('title', '').strip()
        description = request.form.get('description', '').strip()
        level = request.form.get('level', '').strip()
        file = request.files.get('file')
        
        if not title or not description or not level:
            return jsonify({'error': 'Campos requeridos'}), 400
        
        if level not in ['A1', 'A2', 'B1']:
            return jsonify({'error': 'Nivel no válido'}), 400
        
        update_fields = {
            'title': title,
            'description': description,
            'level': level
        }
        
        if file and file.filename != '':
            if not allowed_file(file.filename):
                return jsonify({'error': 'Tipo no permitido'}), 400
            
            upload_result = cloudinary.uploader.upload(
                file,
                folder="clases_ingles",
                resource_type="auto",
                use_filename=True,
                unique_filename=True
            )
            
            update_fields.update({
                'file_url': upload_result['secure_url'],
                'file_name': secure_filename(file.filename),
                'file_type': get_file_type(file.filename),
                'file_size': upload_result.get('bytes', 0)
            })
        
        set_clause = ', '.join([f"{k} = %s" for k in update_fields.keys()])
        values = list(update_fields.values())
        values.append(class_id)
        
        cursor.execute(f'''
            UPDATE classes 
            SET {set_clause}
            WHERE id = %s
        ''', values)
        
        conn.commit()
        cursor.close()
        conn.close()
        
        return jsonify({'message': 'Clase actualizada'}), 200
        
    except Exception as e:
        print(f"❌ Error: {str(e)}")
        return jsonify({'error': 'Error'}), 500

@app.route('/api/classes/<int:class_id>', methods=['DELETE'])
@token_required
@teacher_required
def delete_class(class_id: int):
    try:
        conn = get_db_connection()
        cursor = conn.cursor()
        
        cursor.execute('UPDATE classes SET is_active = FALSE WHERE id = %s', (class_id,))
        conn.commit()
        
        cursor.close()
        conn.close()
        
        return jsonify({'message': 'Clase eliminada'}), 200
        
    except Exception as e:
        print(f"❌ Error: {str(e)}")
        return jsonify({'error': 'Error'}), 500

@app.route('/api/classes/<int:class_id>/download', methods=['GET'])
@token_required
def download_class(class_id: int):
    try:
        conn = get_db_connection()
        cursor = conn.cursor()
        
        if request.user_role == 'student':
            cursor.execute('''
                SELECT sa.has_access FROM student_access sa 
                WHERE sa.student_id = %s AND sa.class_id = %s AND sa.has_access = TRUE
            ''', (request.user_id, class_id))
            
            if not cursor.fetchone():
                cursor.close()
                conn.close()
                return jsonify({'error': 'Sin acceso'}), 403
        
        cursor.execute('SELECT * FROM classes WHERE id = %s AND is_active = TRUE', (class_id,))
        class_data = cursor.fetchone()
        
        if not class_data:
            cursor.close()
            conn.close()
            return jsonify({'error': 'Clase no encontrada'}), 404
        
        cursor.execute('INSERT INTO downloads (user_id, class_id) VALUES (%s, %s)', 
                      (request.user_id, class_id))
        
        if request.user_role == 'student':
            cursor.execute('''
                UPDATE student_access 
                SET downloaded = TRUE, downloaded_at = CURRENT_TIMESTAMP
                WHERE student_id = %s AND class_id = %s
            ''', (request.user_id, class_id))
        
        conn.commit()
        cursor.close()
        conn.close()
        
        return jsonify({
            'download_url': class_data['file_url']
        }), 200
        
    except Exception as e:
        print(f"❌ Error: {str(e)}")
        return jsonify({'error': 'Error'}), 500

# ============ ESTUDIANTES ============
@app.route('/api/students', methods=['GET'])
@token_required
@teacher_required
def get_students():
    try:
        conn = get_db_connection()
        cursor = conn.cursor()
        
        cursor.execute('''
            SELECT u.id, u.username, u.name, u.level, u.avatar_url, u.avatar_version, u.registration_date
            FROM users u 
            WHERE u.role = 'student' AND u.is_active = TRUE
            ORDER BY u.name
        ''')
        
        students = []
        for row in cursor.fetchall():
            students.append(dict(row))
        
        cursor.close()
        conn.close()
        
        return jsonify(students), 200
        
    except Exception as e:
        print(f"❌ Error: {str(e)}")
        return jsonify({'error': 'Error'}), 500

@app.route('/api/students/<int:student_id>/level', methods=['PUT'])
@token_required
@teacher_required
def update_student_level(student_id: int):
    try:
        data = request.get_json()
        
        if not data or 'level' not in data:
            return jsonify({'error': 'Nivel requerido'}), 400
        
        level = data['level'].strip()
        
        conn = get_db_connection()
        cursor = conn.cursor()
        
        cursor.execute('''
            UPDATE users 
            SET level = %s 
            WHERE id = %s AND role = 'student'
            RETURNING id, username, name, level
        ''', (level, student_id))
        
        updated = cursor.fetchone()
        conn.commit()
        cursor.close()
        conn.close()
        
        if not updated:
            return jsonify({'error': 'Estudiante no encontrado'}), 404
        
        socketio.emit('student_level_updated', {
            'student_id': student_id,
            'level': level
        })
        
        return jsonify({
            'message': 'Nivel actualizado',
            'student': dict(updated)
        }), 200
        
    except Exception as e:
        print(f"❌ Error: {str(e)}")
        return jsonify({'error': 'Error'}), 500

@app.route('/api/students/<int:student_id>/access', methods=['GET'])
@token_required
@teacher_required
def get_student_access(student_id: int):
    try:
        conn = get_db_connection()
        cursor = conn.cursor()
        
        cursor.execute('SELECT level FROM users WHERE id = %s', (student_id,))
        student = cursor.fetchone()
        
        if not student:
            cursor.close()
            conn.close()
            return jsonify({'error': 'Estudiante no encontrado'}), 404
        
        student_level = student['level']
        
        cursor.execute('''
            SELECT c.*, 
                   COALESCE(sa.has_access, FALSE) as has_access,
                   sa.downloaded
            FROM classes c
            LEFT JOIN student_access sa ON c.id = sa.class_id AND sa.student_id = %s
            WHERE c.level = %s AND c.is_active = TRUE
            ORDER BY c.created_at DESC
        ''', (student_id, student_level))
        
        classes = []
        for row in cursor.fetchall():
            classes.append(dict(row))
        
        cursor.close()
        conn.close()
        
        return jsonify(classes), 200
        
    except Exception as e:
        print(f"❌ Error: {str(e)}")
        return jsonify({'error': 'Error'}), 500

@app.route('/api/students/<int:student_id>/access', methods=['POST'])
@token_required
@teacher_required
def update_student_access(student_id: int):
    try:
        data = request.get_json()
        
        if not data or 'class_id' not in data or 'has_access' not in data:
            return jsonify({'error': 'Datos incompletos'}), 400
        
        class_id = data['class_id']
        has_access = data['has_access']
        
        conn = get_db_connection()
        cursor = conn.cursor()
        
        cursor.execute('''
            INSERT INTO student_access (student_id, class_id, has_access, granted_by)
            VALUES (%s, %s, %s, %s)
            ON CONFLICT (student_id, class_id) 
            DO UPDATE SET has_access = %s, granted_by = %s, granted_at = CURRENT_TIMESTAMP
        ''', (student_id, class_id, has_access, request.user_id,
              has_access, request.user_id))
        
        conn.commit()
        
        cursor.execute('SELECT title FROM classes WHERE id = %s', (class_id,))
        class_info = cursor.fetchone()
        
        cursor.close()
        conn.close()
        
        socketio.emit('access_updated', {
            'student_id': student_id,
            'class_id': class_id,
            'has_access': has_access,
            'class_title': class_info['title'] if class_info else 'Clase'
        })
        
        return jsonify({'message': 'Acceso actualizado'}), 200
        
    except Exception as e:
        print(f"❌ Error: {str(e)}")
        return jsonify({'error': 'Error'}), 500

# ============ PERFIL ============
@app.route('/api/profile', methods=['GET'])
@token_required
def get_profile():
    try:
        conn = get_db_connection()
        cursor = conn.cursor()
        
        cursor.execute('''
            SELECT id, username, name, role, level, avatar_url, avatar_version, registration_date
            FROM users WHERE id = %s
        ''', (request.user_id,))
        
        user = cursor.fetchone()
        
        if not user:
            cursor.close()
            conn.close()
            return jsonify({'error': 'Usuario no encontrado'}), 404
        
        profile = dict(user)
        
        cursor.close()
        conn.close()
        
        return jsonify(profile), 200
        
    except Exception as e:
        print(f"❌ Error: {str(e)}")
        return jsonify({'error': 'Error'}), 500

# ============ DEEPSEEK CHAT CON STREAMING ============
@app.route('/api/chat/stream', methods=['POST'])
@token_required
def chat_stream():
    """Chat con DeepSeek en tiempo real (streaming)"""
    try:
        data = request.get_json()
        messages = data.get('messages', [])
        temperature = data.get('temperature', 0.7)
        enable_search = data.get('enable_search', False)  # Búsqueda en internet
        
        # Añadir instrucciones del sistema según el rol
        system_message = {
            'role': 'system',
            'content': f"""Eres un asistente experto en enseñanza de inglés. 
                        El usuario es un {currentUser.role} de la plataforma EnglishClasses.
                        Responde de manera clara, educativa y amigable.
                        Puedes ayudar con gramática, vocabulario, pronunciación, ejercicios, etc.
                        {"Puedes buscar información actualizada en internet cuando sea necesario." if enable_search else ""}
                        Si el usuario sube una imagen, extrae el texto y analízalo.
                        Si sube documentos, léelos y responde sobre su contenido."""
        }
        
        full_messages = [system_message] + messages
        
        # Parámetros para la API
        params = {
            'model': 'deepseek-chat',
            'messages': full_messages,
            'temperature': temperature,
            'max_tokens': 4000,
            'stream': True
        }
        
        # Activar búsqueda en internet si se solicita
        if enable_search:
            params['tools'] = [{
                'type': 'web_search',
                'function': {
                    'description': 'Buscar información actualizada en internet'
                }
            }]
        
        # Llamar a DeepSeek API en modo streaming
        response = deepseek_client.chat.completions.create(**params)
        
        def generate():
            for chunk in response:
                if chunk.choices[0].delta.content:
                    yield f"data: {json.dumps({'content': chunk.choices[0].delta.content})}\n\n"
            yield "data: [DONE]\n\n"
        
        return Response(
            stream_with_context(generate()),
            mimetype='text/event-stream',
            headers={
                'Cache-Control': 'no-cache',
                'X-Accel-Buffering': 'no'
            }
        )
        
    except Exception as e:
        print(f"❌ Error en chat stream: {str(e)}")
        return jsonify({'error': str(e)}), 500

# ============ DEEPSEEK CHAT CON ARCHIVOS ============
@app.route('/api/chat-with-files', methods=['POST'])
@token_required
def chat_with_files():
    """Chat con soporte para imágenes y documentos (streaming)"""
    try:
        messages_json = request.form.get('messages', '[]')
        messages = json.loads(messages_json)
        temperature = float(request.form.get('temperature', 0.7))
        enable_search = request.form.get('enable_search', 'false') == 'true'
        
        # Procesar imágenes
        images = request.files.getlist('images')
        image_contents = []
        
        for img in images:
            if img and img.filename:
                img_data = img.read()
                img_base64 = base64.b64encode(img_data).decode('utf-8')
                mime_type = img.mimetype or 'image/jpeg'
                
                image_contents.append({
                    "type": "image_url",
                    "image_url": {
                        "url": f"data:{mime_type};base64,{img_base64}"
                    }
                })
        
        # Procesar documentos
        documents = request.files.getlist('documents')
        document_texts = []
        
        for doc in documents:
            if doc and doc.filename:
                # Guardar posición actual
                pos = doc.tell()
                doc.seek(0)
                
                text = extract_text_from_file(doc, doc.filename)
                document_texts.append(f"--- Documento: {doc.filename} ---\n{text}\n")
                
                # Restaurar posición
                doc.seek(pos)
        
        # Construir mensaje del usuario
        user_content = []
        
        if messages and messages[-1]['role'] == 'user':
            last_message = messages[-1]['content']
            if last_message:
                user_content.append({"type": "text", "text": last_message})
        
        user_content.extend(image_contents)
        
        if document_texts:
            doc_text = "\n\n".join(document_texts)
            user_content.append({"type": "text", "text": f"Contenido de documentos adjuntos:\n{doc_text}"})
        
        # Preparar mensajes completos
        system_message = {
            'role': 'system',
            'content': f"""Eres un asistente experto en enseñanza de inglés. 
                        El usuario es un {currentUser.role} de la plataforma EnglishClasses.
                        Responde de manera clara, educativa y amigable.
                        {"Puedes buscar información actualizada en internet cuando sea necesario." if enable_search else ""}
                        Si hay imágenes, extrae el texto y analízalo.
                        Si hay documentos, úsalos para responder."""
        }
        
        full_messages = [system_message] + messages[:-1] + [{
            'role': 'user',
            'content': user_content
        }]
        
        # Parámetros para la API
        params = {
            'model': 'deepseek-chat',
            'messages': full_messages,
            'temperature': temperature,
            'max_tokens': 4000,
            'stream': True
        }
        
        if enable_search:
            params['tools'] = [{
                'type': 'web_search',
                'function': {
                    'description': 'Buscar información actualizada en internet'
                }
            }]
        
        # Llamar a DeepSeek API
        response = deepseek_client.chat.completions.create(**params)
        
        def generate():
            for chunk in response:
                if chunk.choices[0].delta.content:
                    yield f"data: {json.dumps({'content': chunk.choices[0].delta.content})}\n\n"
            yield "data: [DONE]\n\n"
        
        return Response(
            stream_with_context(generate()),
            mimetype='text/event-stream',
            headers={
                'Cache-Control': 'no-cache',
                'X-Accel-Buffering': 'no'
            }
        )
        
    except Exception as e:
        print(f"❌ Error en chat con archivos: {str(e)}")
        return jsonify({'error': str(e)}), 500

# ============ WEBSOCKET ============
@socketio.on('connect')
def handle_connect():
    emit('connected', {'message': 'Conectado al servidor'})

# ============ ARCHIVOS ESTÁTICOS ============
@app.route('/<path:path>')
def serve_static(path):
    try:
        return send_from_directory('.', path)
    except:
        return jsonify({'error': 'Archivo no encontrado'}), 404

# ============ INICIALIZACIÓN ============
if __name__ == '__main__':
    print("=" * 60)
    print("🚀 SERVIDOR DE CLASES DE INGLÉS - V7.0")
    print("=" * 60)
    
    # Inicializar base de datos
    init_db()
    
    print(f"🔑 Código estudiante: {STUDENT_CODE}")
    print(f"👨‍🏫 Código profesor: {TEACHER_CODE}")
    print("=" * 60)
    
    port = int(os.environ.get('PORT', 5000))
    
    socketio.run(
        app, 
        host='0.0.0.0', 
        port=port, 
        debug=False,
        allow_unsafe_werkzeug=True
    )